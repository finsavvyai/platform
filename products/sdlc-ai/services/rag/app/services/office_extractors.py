"""
Microsoft Office document extractors for SDLC.ai platform.

This module provides comprehensive extraction capabilities for DOCX, XLSX, and PPTX formats
with support for embedded content, metadata preservation, and structure analysis.
"""

import io
import logging
import re
import zipfile
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple, Union

import pandas as pd
from docx import Document as DocxDocument
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Inches
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from app.services.document_processor import BaseExtractor, ExtractionResult, ProcessingOptions

logger = logging.getLogger(__name__)


class DOCXExtractor(BaseExtractor):
    """High-precision Microsoft Word document extractor."""

    def supports_format(self, content_type: str) -> bool:
        """Check if DOCX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ]

    async def extract(self, file_data: bytes, options: ProcessingOptions) -> ExtractionResult:
        """Extract text and structure from DOCX documents."""
        start_time = datetime.now()

        try:
            document = DocxDocument(io.BytesIO(file_data))

            # Extract different content types
            main_text = self._extract_main_text(document)
            headers_footers = self._extract_headers_footers(document)
            tables = self._extract_tables(document)
            images = self._extract_images(document)
            metadata = self._extract_docx_metadata(document)
            structure = self._analyze_docx_structure(document)

            # Combine all text
            full_text = self._combine_text_content(main_text, headers_footers, tables)

            # Calculate quality metrics
            quality_metrics = self._calculate_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=[{"content": main_text, "type": "main"}],  # DOCX doesn't have pages like PDF
                tables=tables,
                images=images,
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.95  # DOCX extraction is generally high quality
            )

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise

    def _extract_main_text(self, document: DocxDocument) -> List[Dict[str, Any]]:
        """Extract main text content with structure information."""
        content = []

        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                para_info = {
                    "text": paragraph.text,
                    "style": paragraph.style.name if paragraph.style else "Normal",
                    "alignment": self._get_alignment_name(paragraph.alignment),
                    "level": self._get_paragraph_level(paragraph),
                    "is_heading": self._is_heading(paragraph),
                    "position": len(content)
                }
                content.append(para_info)

        return content

    def _extract_headers_footers(self, document: DocxDocument) -> Dict[str, List[str]]:
        """Extract headers and footers from all sections."""
        headers = []
        footers = []

        for section in document.sections:
            # Header
            if section.header:
                for paragraph in section.header.paragraphs:
                    if paragraph.text.strip():
                        headers.append(paragraph.text)

            # Footer
            if section.footer:
                for paragraph in section.footer.paragraphs:
                    if paragraph.text.strip():
                        footers.append(paragraph.text)

        return {
            "headers": headers,
            "footers": footers
        }

    def _extract_tables(self, document: DocxDocument) -> List[Dict[str, Any]]:
        """Extract tables with structure and formatting information."""
        tables = []

        for table_idx, table in enumerate(document.tables):
            table_data = {
                "table_index": table_idx,
                "rows": [],
                "row_count": len(table.rows),
                "col_count": len(table.columns) if table.rows else 0,
                "has_header": self._detect_table_header(table),
                "formatting": self._analyze_table_formatting(table)
            }

            for row_idx, row in enumerate(table.rows):
                row_data = {
                    "row_index": row_idx,
                    "cells": [],
                    "is_header": row_idx == 0 and table_data["has_header"]
                }

                for cell_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    cell_info = {
                        "column_index": cell_idx,
                        "text": cell_text,
                        "colspan": getattr(cell, '_element', {}).get('gridSpan', 1),
                        "rowspan": getattr(cell, '_element', {}).get('vMerge', 1),
                        "is_empty": not cell_text
                    }
                    row_data["cells"].append(cell_info)

                table_data["rows"].append(row_data)

            tables.append(table_data)

        return tables

    def _extract_images(self, document: DocxDocument) -> List[Dict[str, Any]]:
        """Extract image information from the document."""
        images = []

        # Extract inline shapes (images)
        for shape in document.inline_shapes:
            if shape.type == 3:  # Inline picture
                image_info = {
                    "type": "inline_image",
                    "width": shape.width,
                    "height": shape.height,
                    "format": "unknown"  # DOCX doesn't easily expose image format
                }
                images.append(image_info)

        return images

    def _extract_docx_metadata(self, document: DocxDocument) -> Dict[str, Any]:
        """Extract document metadata."""
        metadata = {}

        # Core properties
        if document.core_properties:
            metadata.update({
                "title": document.core_properties.title or "",
                "author": document.core_properties.author or "",
                "subject": document.core_properties.subject or "",
                "keywords": document.core_properties.keywords or "",
                "category": document.core_properties.category or "",
                "comments": document.core_properties.comments or "",
                "created": document.core_properties.created,
                "modified": document.core_properties.modified,
                "last_modified_by": document.core_properties.last_modified_by or "",
                "revision": document.core_properties.revision or 0,
            })

        # Document statistics
        metadata.update({
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
            "section_count": len(document.sections),
            "inline_shape_count": len(document.inline_shapes),
        })

        return metadata

    def _analyze_docx_structure(self, document: DocxDocument) -> Dict[str, Any]:
        """Analyze document structure."""
        structure = {
            "has_toc": False,
            "sections": [],
            "headings": [],
            "styles_used": set(),
            "layout_complexity": "simple"
        }

        # Analyze paragraphs for structure
        for para_idx, paragraph in enumerate(document.paragraphs):
            if paragraph.style:
                structure["styles_used"].add(paragraph.style.name)

            # Detect headings
            if self._is_heading(paragraph):
                heading_info = {
                    "text": paragraph.text,
                    "level": self._get_paragraph_level(paragraph),
                    "style": paragraph.style.name,
                    "position": para_idx
                }
                structure["headings"].append(heading_info)

            # Detect table of contents
            if any(keyword in paragraph.text.lower() for keyword in ["contents", "table of contents", "index"]):
                structure["has_toc"] = True

        # Convert set to list for JSON serialization
        structure["styles_used"] = list(structure["styles_used"])

        # Determine layout complexity
        if len(document.tables) > 5 or len(structure["headings"]) > 10:
            structure["layout_complexity"] = "complex"
        elif len(document.tables) > 2 or len(structure["headings"]) > 5:
            structure["layout_complexity"] = "medium"

        # Analyze sections
        for section_idx, section in enumerate(document.sections):
            section_info = {
                "section_index": section_idx,
                "has_header": bool(section.header),
                "has_footer": bool(section.footer),
                "page_orientation": "portrait" if section.orientation == 0 else "landscape"
            }
            structure["sections"].append(section_info)

        return structure

    def _get_alignment_name(self, alignment) -> str:
        """Convert paragraph alignment to readable name."""
        alignment_map = {
            WD_PARAGRAPH_ALIGNMENT.LEFT: "left",
            WD_PARAGRAPH_ALIGNMENT.CENTER: "center",
            WD_PARAGRAPH_ALIGNMENT.RIGHT: "right",
            WD_PARAGRAPH_ALIGNMENT.JUSTIFY: "justify",
            None: "left"
        }
        return alignment_map.get(alignment, "left")

    def _get_paragraph_level(self, paragraph) -> int:
        """Get the hierarchical level of a paragraph."""
        style_name = paragraph.style.name.lower() if paragraph.style else ""

        if "heading" in style_name:
            # Extract number from heading style
            match = re.search(r'heading\s*(\d+)', style_name)
            if match:
                return int(match.group(1))
        elif "title" in style_name:
            return 1

        return 10  # Regular text has low level

    def _is_heading(self, paragraph) -> bool:
        """Check if paragraph is a heading."""
        style_name = paragraph.style.name.lower() if paragraph.style else ""
        return "heading" in style_name or "title" in style_name

    def _detect_table_header(self, table) -> bool:
        """Detect if table has a header row."""
        if not table.rows:
            return False

        first_row = table.rows[0]
        # Check if first row has different formatting or contains header-like text
        header_indicators = [
            any(cell.text.strip().isupper() for cell in first_row.cells if cell.text.strip()),
            len(first_row.cells) > 1,  # Multiple columns suggest structured data
            any(keyword in " ".join(cell.text.lower() for cell in first_row.cells if cell.text.strip())
                for keyword in ["name", "date", "id", "type", "description", "status"])
        ]

        return sum(header_indicators) >= 2

    def _analyze_table_formatting(self, table) -> Dict[str, Any]:
        """Analyze table formatting."""
        return {
            "has_borders": True,  # DOCX tables typically have borders
            "row_count": len(table.rows),
            "col_count": len(table.columns) if table.rows else 0,
            "is_complex": len(table.rows) > 10 or (table.rows and len(table.columns) > 5)
        }

    def _combine_text_content(self, main_text: List[Dict], headers_footers: Dict, tables: List[Dict]) -> str:
        """Combine all text content into a single string."""
        text_parts = []

        # Add headers
        for header in headers_footers.get("headers", []):
            text_parts.append(f"[HEADER] {header}")

        text_parts.append("")  # Add separator

        # Add main text
        for para in main_text:
            if para["is_heading"]:
                text_parts.append(f"\n{para['text']}\n")
            else:
                text_parts.append(para["text"])

        # Add tables
        for table in tables:
            text_parts.append("\n[TABLE]")
            for row in table["rows"]:
                row_text = " | ".join(cell["text"] for cell in row["cells"])
                text_parts.append(row_text)
            text_parts.append("[/TABLE]\n")

        # Add footers
        for footer in headers_footers.get("footers", []):
            text_parts.append(f"[FOOTER] {footer}")

        return "\n".join(text_parts)

    def _calculate_quality_metrics(self, text: str, metadata: Dict[str, Any]) -> Dict[str, float]:
        """Calculate quality metrics for DOCX extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "word_count": len(text.split()),
            "paragraph_count": metadata.get("paragraph_count", 0),
            "table_count": metadata.get("table_count", 0),
            "has_structure": bool(metadata.get("headings", [])),
            "completeness_score": 0.0,
        }

        # Calculate completeness score
        completeness_factors = [
            len(text) > 500,  # Substantial content
            metrics["paragraph_count"] > 5,  # Multiple paragraphs
            metadata.get("title", "") != "",  # Has title
            metrics["has_structure"],  # Has document structure
        ]

        metrics["completeness_score"] = sum(completeness_factors) / len(completeness_factors)

        # Overall quality
        metrics["overall_quality"] = (
            min(1.0, len(text) / 5000) * 0.3 +  # Content length
            min(1.0, metrics["paragraph_count"] / 50) * 0.3 +  # Paragraph count
            metrics["completeness_score"] * 0.4  # Completeness
        )

        return metrics


class XLSXExtractor(BaseExtractor):
    """Microsoft Excel document extractor."""

    def supports_format(self, content_type: str) -> bool:
        """Check if XLSX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ]

    async def extract(self, file_data: bytes, options: ProcessingOptions) -> ExtractionResult:
        """Extract text and data from Excel documents."""
        start_time = datetime.now()

        try:
            workbook = load_workbook(io.BytesIO(file_data), data_only=True)

            # Extract data from all sheets
            sheets_data = []
            tables = []
            metadata = self._extract_xlsx_metadata(workbook)

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = self._extract_sheet_data(sheet, sheet_name)
                sheets_data.append(sheet_data)

                # Extract tables from the sheet
                sheet_tables = self._extract_sheet_tables(sheet, sheet_name)
                tables.extend(sheet_tables)

            # Combine all text
            full_text = self._combine_sheet_text(sheets_data)

            # Create structure information
            structure = {
                "sheet_count": len(workbook.sheetnames),
                "sheets": [{"name": sheet_name, "has_data": True} for sheet_name in workbook.sheetnames],
                "total_tables": len(tables)
            }

            # Calculate quality metrics
            quality_metrics = self._calculate_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=[],  # Excel doesn't have pages
                tables=tables,
                images=[],
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.90
            )

        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise

    def _extract_sheet_data(self, sheet: Worksheet, sheet_name: str) -> Dict[str, Any]:
        """Extract data from a specific sheet."""
        sheet_data = {
            "sheet_name": sheet_name,
            "rows": [],
            "used_range": None,
            "has_data": False,
            "data_types": set()
        }

        # Find the used range
        if sheet.max_row > 1 and sheet.max_column > 1:
            sheet_data["used_range"] = {
                "min_row": sheet.min_row,
                "max_row": sheet.max_row,
                "min_col": sheet.min_column,
                "max_col": sheet.max_column
            }

            # Extract data row by row
            for row_idx in range(sheet.min_row, sheet.max_row + 1):
                row_data = {
                    "row_index": row_idx,
                    "cells": [],
                    "has_data": False
                }

                for col_idx in range(sheet.min_column, sheet.max_column + 1):
                    cell = sheet.cell(row=row_idx, column=col_idx)
                    cell_value = self._get_cell_value(cell)

                    cell_info = {
                        "column_letter": get_column_letter(col_idx),
                        "column_index": col_idx,
                        "row_index": row_idx,
                        "value": cell_value,
                        "data_type": type(cell.value).__name__ if cell.value is not None else "empty",
                        "is_empty": cell.value is None or str(cell.value).strip() == "",
                        "has_formula": bool(cell.data_type == 'f')  # Formula cell
                    }

                    if not cell_info["is_empty"]:
                        row_data["has_data"] = True
                        sheet_data["has_data"] = True
                        sheet_data["data_types"].add(cell_info["data_type"])

                    row_data["cells"].append(cell_info)

                sheet_data["rows"].append(row_data)

        # Convert set to list for serialization
        sheet_data["data_types"] = list(sheet_data["data_types"])

        return sheet_data

    def _get_cell_value(self, cell) -> str:
        """Get formatted string value from a cell."""
        if cell.value is None:
            return ""
        elif isinstance(cell.value, (int, float)):
            return str(cell.value)
        elif isinstance(cell.value, datetime):
            return cell.value.strftime("%Y-%m-%d %H:%M:%S")
        else:
            return str(cell.value)

    def _extract_sheet_tables(self, sheet: Worksheet, sheet_name: str) -> List[Dict[str, Any]]:
        """Extract structured tables from a sheet."""
        tables = []

        # Look for Excel tables (if they exist)
        if hasattr(sheet, 'tables') and sheet.tables:
            for table in sheet.tables:
                table_data = {
                    "sheet_name": sheet_name,
                    "table_name": table.name,
                    "range": str(table.ref),
                    "has_header": table.headerRowCount > 0,
                    "rows": []
                }

                # Extract table data
                for row in sheet[table.ref]:
                    row_data = [self._get_cell_value(cell) for cell in row]
                    table_data["rows"].append(row_data)

                tables.append(table_data)

        # If no formal tables, try to detect data tables
        if not tables:
            sheet_data_rows = self._detect_data_tables(sheet)
            if sheet_data_rows:
                tables.extend(sheet_data_rows)

        return tables

    def _detect_data_tables(self, sheet: Worksheet) -> List[Dict[str, Any]]:
        """Detect and extract informal data tables."""
        tables = []

        # Simple table detection: look for rectangular data blocks
        if sheet.max_row <= 1 or sheet.max_column <= 1:
            return tables

        # Find contiguous data ranges
        data_ranges = []
        current_start_row = None
        current_end_row = None

        for row_idx in range(sheet.min_row, sheet.max_row + 1):
            row_has_data = any(sheet.cell(row=row_idx, column=col_idx).value is not None
                             for col_idx in range(sheet.min_column, sheet.max_column + 1))

            if row_has_data and current_start_row is None:
                current_start_row = row_idx
            elif row_has_data and current_start_row is not None:
                current_end_row = row_idx
            elif not row_has_data and current_start_row is not None:
                data_ranges.append((current_start_row, current_end_row))
                current_start_row = None
                current_end_row = None

        # Add the last range if the sheet ends with data
        if current_start_row is not None:
            data_ranges.append((current_start_row, current_end_row or sheet.max_row))

        # Convert ranges to table data
        for range_start, range_end in data_ranges:
            if range_end - range_start >= 2:  # At least 3 rows to consider it a table
                table_data = {
                    "sheet_name": sheet.title,
                    "table_name": f"detected_table_{len(tables) + 1}",
                    "range": f"A{range_start}:{get_column_letter(sheet.max_column)}{range_end}",
                    "has_header": True,  # Assume first row is header
                    "rows": []
                }

                for row_idx in range(range_start, range_end + 1):
                    row_data = []
                    for col_idx in range(sheet.min_column, sheet.max_column + 1):
                        cell = sheet.cell(row=row_idx, column=col_idx)
                        row_data.append(self._get_cell_value(cell))
                    table_data["rows"].append(row_data)

                tables.append(table_data)

        return tables

    def _extract_xlsx_metadata(self, workbook) -> Dict[str, Any]:
        """Extract workbook metadata."""
        metadata = {}

        # Core properties
        if workbook.properties:
            metadata.update({
                "title": workbook.properties.title or "",
                "author": workbook.properties.creator or "",
                "subject": workbook.properties.subject or "",
                "description": workbook.properties.description or "",
                "keywords": workbook.properties.keywords or "",
                "category": workbook.properties.category or "",
                "created": workbook.properties.created,
                "modified": workbook.properties.modified,
                "last_modified_by": workbook.properties.lastModifiedBy or "",
            })

        # Workbook statistics
        metadata.update({
            "sheet_count": len(workbook.sheetnames),
            "sheet_names": workbook.sheetnames,
            "calculation_mode": workbook.calculation.calcMode,
            "has_external_links": len(workbook._external_links) > 0 if hasattr(workbook, '_external_links') else False,
        })

        return metadata

    def _combine_sheet_text(self, sheets_data: List[Dict[str, Any]]) -> str:
        """Combine text from all sheets."""
        text_parts = []

        for sheet in sheets_data:
            if not sheet["has_data"]:
                continue

            text_parts.append(f"\n[SHEET: {sheet['sheet_name']}]")

            for row in sheet["rows"]:
                if row["has_data"]:
                    row_text = " | ".join(cell["value"] for cell in row["cells"] if not cell["is_empty"])
                    if row_text.strip():
                        text_parts.append(row_text)

            text_parts.append("[/SHEET]\n")

        return "\n".join(text_parts)

    def _calculate_quality_metrics(self, text: str, metadata: Dict[str, Any]) -> Dict[str, float]:
        """Calculate quality metrics for Excel extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "sheet_count": metadata.get("sheet_count", 0),
            "data_density": 0.0,
            "structure_score": 0.0,
            "completeness_score": 0.0,
        }

        # Calculate data density (ratio of data cells to potential cells)
        total_cells = sum(len(row["cells"]) for sheet in [] for row in sheet.get("rows", []))
        non_empty_cells = sum(1 for sheet in [] for row in sheet.get("rows", [])
                            for cell in row.get("cells", []) if not cell.get("is_empty", True))

        if total_cells > 0:
            metrics["data_density"] = non_empty_cells / total_cells

        # Structure score based on sheet organization
        metrics["structure_score"] = min(1.0, metadata.get("sheet_count", 0) / 10)

        # Completeness score
        completeness_factors = [
            len(text) > 200,  # Substantial content
            metadata.get("sheet_count", 0) > 1,  # Multiple sheets
            metrics["data_density"] > 0.1,  # Good data density
        ]

        metrics["completeness_score"] = sum(completeness_factors) / len(completeness_factors)

        # Overall quality
        metrics["overall_quality"] = (
            metrics["data_density"] * 0.4 +
            metrics["structure_score"] * 0.3 +
            metrics["completeness_score"] * 0.3
        )

        return metrics


class PPTXExtractor(BaseExtractor):
    """Microsoft PowerPoint document extractor."""

    def supports_format(self, content_type: str) -> bool:
        """Check if PPTX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]

    async def extract(self, file_data: bytes, options: ProcessingOptions) -> ExtractionResult:
        """Extract text and content from PowerPoint presentations."""
        start_time = datetime.now()

        try:
            presentation = Presentation(io.BytesIO(file_data))

            # Extract content from slides
            slides_data = []
            for slide_idx, slide in enumerate(presentation.slides):
                slide_data = self._extract_slide_content(slide, slide_idx)
                slides_data.append(slide_data)

            # Extract presentation metadata
            metadata = self._extract_pptx_metadata(presentation)

            # Combine all text
            full_text = self._combine_slide_text(slides_data)

            # Analyze presentation structure
            structure = self._analyze_pptx_structure(presentation, slides_data)

            # Extract images and media
            images = self._extract_media_info(presentation)

            # Calculate quality metrics
            quality_metrics = self._calculate_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=[{"slide_number": i + 1, "content": slide} for i, slide in enumerate(slides_data)],
                tables=[],  # Tables in PPTX are handled within slide content
                images=images,
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.88
            )

        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise

    def _extract_slide_content(self, slide, slide_idx: int) -> Dict[str, Any]:
        """Extract content from a single slide."""
        slide_data = {
            "slide_index": slide_idx,
            "slide_number": slide_idx + 1,
            "title": "",
            "content": [],
            "notes": "",
            "has_content": False,
            "content_types": set()
        }

        # Extract title (usually the first text box or placeholder)
        if hasattr(slide, 'shapes') and slide.shapes:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Check if this might be a title (usually first or prominent text)
                    if len(shape.text.strip()) < 100 and (not slide_data["title"] or len(shape.text) > len(slide_data["title"])):
                        slide_data["title"] = shape.text.strip()
                    else:
                        slide_data["content"].append(shape.text.strip())
                        slide_data["content_types"].add("text")
                        slide_data["has_content"] = True

        # Extract notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_text = []
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    notes_text.append(shape.text.strip())
            slide_data["notes"] = "\n".join(notes_text)

        # Convert set to list for serialization
        slide_data["content_types"] = list(slide_data["content_types"])

        return slide_data

    def _extract_pptx_metadata(self, presentation) -> Dict[str, Any]:
        """Extract presentation metadata."""
        metadata = {}

        # Core properties
        if hasattr(presentation, 'core_properties') and presentation.core_properties:
            core_props = presentation.core_properties
            metadata.update({
                "title": getattr(core_props, 'title', '') or "",
                "author": getattr(core_props, 'author', '') or "",
                "subject": getattr(core_props, 'subject', '') or "",
                "keywords": getattr(core_props, 'keywords', '') or "",
                "comments": getattr(core_props, 'comments', '') or "",
                "created": getattr(core_props, 'created', None),
                "modified": getattr(core_props, 'modified', None),
                "last_modified_by": getattr(core_props, 'last_modified_by', '') or "",
                "revision": getattr(core_props, 'revision', 0) or 0,
            })

        # Presentation statistics
        metadata.update({
            "slide_count": len(presentation.slides),
            "has_notes": any(hasattr(slide, 'notes_slide') and slide.notes_slide for slide in presentation.slides),
        })

        # Slide layout information
        slide_layouts = set()
        for slide in presentation.slides:
            if hasattr(slide, 'slide_layout') and slide.slide_layout:
                slide_layouts.add(type(slide.slide_layout).__name__)

        metadata["slide_layout_types"] = list(slide_layouts)

        return metadata

    def _analyze_pptx_structure(self, presentation, slides_data: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Analyze presentation structure."""
        structure = {
            "slide_count": len(presentation.slides),
            "slides_with_titles": 0,
            "slides_with_content": 0,
            "slides_with_notes": 0,
            "content_distribution": {},
            "presentation_type": "unknown"
        }

        content_counts = {"title_only": 0, "title_and_content": 0, "content_only": 0, "empty": 0}

        for slide in slides_data:
            # Count slides with different content types
            if slide["title"] and slide["content"]:
                content_counts["title_and_content"] += 1
                structure["slides_with_titles"] += 1
                structure["slides_with_content"] += 1
            elif slide["title"] and not slide["content"]:
                content_counts["title_only"] += 1
                structure["slides_with_titles"] += 1
            elif not slide["title"] and slide["content"]:
                content_counts["content_only"] += 1
                structure["slides_with_content"] += 1
            else:
                content_counts["empty"] += 1

            if slide["notes"]:
                structure["slides_with_notes"] += 1

        structure["content_distribution"] = content_counts

        # Determine presentation type based on content distribution
        title_slides = content_counts["title_only"]
        content_slides = content_counts["title_and_content"] + content_counts["content_only"]

        if title_slides > 0 and content_slides > 0:
            if content_slides / len(slides_data) > 0.7:
                structure["presentation_type"] = "standard_presentation"
            else:
                structure["presentation_type"] = "title_heavy"
        elif content_slides > 0:
            structure["presentation_type"] = "content_focused"
        else:
            structure["presentation_type"] = "minimal"

        return structure

    def _extract_media_info(self, presentation) -> List[Dict[str, Any]]:
        """Extract information about images and media in the presentation."""
        media = []

        for slide_idx, slide in enumerate(presentation.slides):
            if hasattr(slide, 'shapes'):
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_info = {
                        "slide_number": slide_idx + 1,
                        "shape_index": shape_idx,
                        "shape_type": "unknown"
                    }

                    # Determine shape type
                    if hasattr(shape, 'shape_type'):
                        try:
                            shape_type_name = shape.shape_type
                            shape_info["shape_type"] = shape_type_name.name if hasattr(shape_type_name, 'name') else str(shape_type_name)
                        except:
                            pass

                    # Check for images
                    if shape_info["shape_type"] == "PICTURE" or hasattr(shape, 'image'):
                        shape_info.update({
                            "media_type": "image",
                            "has_image": True
                        })
                        media.append(shape_info)

                    # Check for charts
                    elif shape_info["shape_type"] == "CHART":
                        shape_info.update({
                            "media_type": "chart",
                            "has_chart": True
                        })
                        media.append(shape_info)

                    # Check for tables
                    elif shape_info["shape_type"] == "TABLE":
                        shape_info.update({
                            "media_type": "table",
                            "has_table": True
                        })
                        media.append(shape_info)

        return media

    def _combine_slide_text(self, slides_data: List[Dict[str, Any]]) -> str:
        """Combine text from all slides."""
        text_parts = []

        for slide in slides_data:
            # Add slide title
            if slide["title"]:
                text_parts.append(f"\n=== SLIDE {slide['slide_number']}: {slide['title']} ===\n")
            else:
                text_parts.append(f"\n=== SLIDE {slide['slide_number']} ===\n")

            # Add slide content
            for content in slide["content"]:
                if content:
                    text_parts.append(content)

            # Add notes if available
            if slide["notes"]:
                text_parts.append(f"\n[NOTES: {slide['notes']}]")

            text_parts.append("")  # Add blank line between slides

        return "\n".join(text_parts)

    def _calculate_quality_metrics(self, text: str, metadata: Dict[str, Any]) -> Dict[str, float]:
        """Calculate quality metrics for PowerPoint extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "slide_count": metadata.get("slide_count", 0),
            "content_per_slide": 0.0,
            "structure_quality": 0.0,
            "completeness_score": 0.0,
        }

        # Calculate content per slide
        if metrics["slide_count"] > 0:
            metrics["content_per_slide"] = len(text) / metrics["slide_count"]

        # Structure quality based on title/content ratio
        structure = metadata.get("structure", {})
        slides_with_titles = structure.get("slides_with_titles", 0)
        slides_with_content = structure.get("slides_with_content", 0)

        if metrics["slide_count"] > 0:
            metrics["structure_quality"] = (slides_with_titles + slides_with_content) / (2 * metrics["slide_count"])

        # Completeness score
        completeness_factors = [
            len(text) > 200,  # Substantial content
            metrics["slide_count"] > 1,  # Multiple slides
            slides_with_titles > 0,  # Has titles
            metrics["content_per_slide"] > 50,  # Good content per slide
        ]

        metrics["completeness_score"] = sum(completeness_factors) / len(completeness_factors)

        # Overall quality
        metrics["overall_quality"] = (
            min(1.0, metrics["content_per_slide"] / 200) * 0.3 +  # Content density
            metrics["structure_quality"] * 0.4 +  # Structure
            metrics["completeness_score"] * 0.3  # Completeness
        )

        return metrics
