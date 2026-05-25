"""
Microsoft Office document extractors for SDLC.ai platform.

This module provides comprehensive text extraction from DOCX, XLSX, and PPTX formats
with metadata preservation and structure analysis.
"""

import io
import logging
import re
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd
from docx import Document as DocxDocument
from docx.opc.exceptions import PackageNotFoundError
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.opc.exceptions import PackageNotFoundError

from ..document_processor import BaseExtractor, ExtractionResult, ProcessingOptions

logger = logging.getLogger(__name__)


class DOCXExtractor(BaseExtractor):
    """High-precision DOCX text extraction with structure analysis."""

    def supports_format(self, content_type: str) -> bool:
        """Check if DOCX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ]

    async def extract(
        self, file_data: bytes, options: ProcessingOptions
    ) -> ExtractionResult:
        """Extract text and structure from DOCX document."""
        start_time = datetime.now()

        try:
            doc = DocxDocument(io.BytesIO(file_data))

            # Extract different content types
            content_parts = []
            metadata = {}
            structure = {}
            tables = []
            images = []
            pages = []

            # Extract main paragraphs with structure
            paragraphs_data = self._extract_paragraphs(doc)
            content_parts.append(paragraphs_data["text"])
            structure.update(paragraphs_data["structure"])

            # Extract tables
            doc_tables = self._extract_tables(doc)
            tables.extend(doc_tables)
            if doc_tables:
                content_parts.append("\n\n".join(table["text"] for table in doc_tables))

            # Extract headers and footers
            headers_footers = self._extract_headers_footers(doc)
            if headers_footers:
                content_parts.append(headers_footers)

            # Extract document properties
            metadata = self._extract_docx_metadata(doc)

            # Combine all content
            full_text = "\n\n".join(filter(None, content_parts))

            # Create page information (estimated for DOCX)
            pages = self._create_docx_page_info(doc, full_text)

            # Analyze document structure
            structure.update(self._analyze_docx_structure(doc))

            # Calculate quality metrics
            quality_metrics = self._calculate_docx_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=pages,
                tables=tables,
                images=images,
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.95,
            )

        except PackageNotFoundError:
            logger.error("Invalid DOCX file format")
            raise ValueError("Invalid DOCX file format")
        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise

    def _extract_paragraphs(self, doc: DocxDocument) -> Dict[str, Any]:
        """Extract paragraphs with style information."""
        paragraphs = []
        structure = {"headings": [], "paragraphs": [], "lists": [], "styles": set()}

        for i, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue

            # Get style information
            style_name = paragraph.style.name if paragraph.style else "Normal"
            structure["styles"].add(style_name)

            # Check if it's a heading
            if style_name.lower().startswith("heading"):
                structure["headings"].append(
                    {
                        "index": i,
                        "text": text,
                        "level": self._extract_heading_level(style_name),
                        "style": style_name,
                    }
                )

            # Check if it's a list item
            if self._is_list_item(paragraph):
                structure["lists"].append(
                    {
                        "index": i,
                        "text": text,
                        "level": self._get_list_level(paragraph),
                        "style": style_name,
                    }
                )

            structure["paragraphs"].append(
                {"index": i, "text": text, "style": style_name, "length": len(text)}
            )

            paragraphs.append(text)

        structure["styles"] = list(structure["styles"])

        return {"text": "\n".join(paragraphs), "structure": structure}

    def _extract_tables(self, doc: DocxDocument) -> List[Dict[str, Any]]:
        """Extract tables from DOCX document."""
        tables = []

        for table_idx, table in enumerate(doc.tables):
            table_text = []
            table_data = []

            for row_idx, row in enumerate(table.rows):
                row_data = []
                row_text = []

                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_data.append(cell_text)
                    row_text.append(cell_text)

                table_data.append(row_data)
                table_text.append(" | ".join(row_text))

            if table_data and len(table_data) > 1:  # At least header + 1 row
                tables.append(
                    {
                        "table_index": table_idx,
                        "rows": table_data,
                        "row_count": len(table_data),
                        "col_count": len(table_data[0]) if table_data else 0,
                        "text": "\n".join(table_text),
                        "confidence": 0.95,
                    }
                )

        return tables

    def _extract_headers_footers(self, doc: DocxDocument) -> Optional[str]:
        """Extract headers and footers from DOCX."""
        header_footer_parts = []

        # Extract headers
        for section in doc.sections:
            if section.header:
                for paragraph in section.header.paragraphs:
                    if paragraph.text.strip():
                        header_footer_parts.append(f"HEADER: {paragraph.text.strip()}")

            if section.footer:
                for paragraph in section.footer.paragraphs:
                    if paragraph.text.strip():
                        header_footer_parts.append(f"FOOTER: {paragraph.text.strip()}")

        return "\n".join(header_footer_parts) if header_footer_parts else None

    def _extract_docx_metadata(self, doc: DocxDocument) -> Dict[str, Any]:
        """Extract metadata from DOCX document properties."""
        metadata = {}

        try:
            core_props = doc.core_properties
            metadata.update(
                {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "keywords": core_props.keywords or "",
                    "category": core_props.category or "",
                    "comments": core_props.comments or "",
                    "created": core_props.created.isoformat()
                    if core_props.created
                    else None,
                    "modified": core_props.modified.isoformat()
                    if core_props.modified
                    else None,
                    "last_modified_by": core_props.last_modified_by or "",
                    "revision": str(core_props.revision)
                    if core_props.revision
                    else "0",
                }
            )
        except Exception as e:
            logger.warning(f"Failed to extract document properties: {e}")

        # Add document statistics
        metadata.update(
            {
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "section_count": len(doc.sections),
            }
        )

        return metadata

    def _extract_heading_level(self, style_name: str) -> int:
        """Extract heading level from style name."""
        match = re.search(r"heading\s*(\d+)", style_name.lower())
        return int(match.group(1)) if match else 1

    def _is_list_item(self, paragraph) -> bool:
        """Check if paragraph is a list item."""
        style_name = paragraph.style.name.lower() if paragraph.style else ""
        return any(list_type in style_name for list_type in ["list", "bullet"])

    def _get_list_level(self, paragraph) -> int:
        """Get list level for paragraph."""
        if paragraph.style and hasattr(paragraph.style, "indent"):
            return getattr(paragraph.style, "indent", 0)
        return 0

    def _analyze_docx_structure(self, doc: DocxDocument) -> Dict[str, Any]:
        """Analyze DOCX document structure."""
        structure = {
            "has_title": False,
            "has_abstract": False,
            "has_toc": False,
            "sections": [],
            "document_type": "unknown",
        }

        # Analyze document content for structure
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        # Check for title
        if doc.paragraphs:
            first_para = doc.paragraphs[0]
            if (first_para.style and "title" in first_para.style.name.lower()) or len(
                first_para.text
            ) < 100:
                structure["has_title"] = True

        # Check for abstract
        abstract_patterns = [r"abstract", r"summary", r"overview"]
        structure["has_abstract"] = any(
            re.search(pattern, full_text, re.IGNORECASE)
            for pattern in abstract_patterns
        )

        # Check for table of contents
        toc_patterns = [r"contents?", r"table of contents", r"index"]
        structure["has_toc"] = any(
            re.search(pattern, full_text, re.IGNORECASE) for pattern in toc_patterns
        )

        # Detect document type
        if any(
            keyword in full_text.lower()
            for keyword in ["chapter", "section", "appendix"]
        ):
            structure["document_type"] = "book"
        elif any(
            keyword in full_text.lower()
            for keyword in ["introduction", "methodology", "conclusion"]
        ):
            structure["document_type"] = "academic"
        elif any(
            keyword in full_text.lower()
            for keyword in ["executive", "summary", "recommendation"]
        ):
            structure["document_type"] = "business"

        return structure

    def _create_docx_page_info(
        self, doc: DocxDocument, text: str
    ) -> List[Dict[str, Any]]:
        """Create estimated page information for DOCX."""
        # Estimate pages based on content (DOCX doesn't have fixed pages)
        words = text.split()
        estimated_pages = max(1, len(words) // 250)  # ~250 words per page

        pages = []
        for i in range(estimated_pages):
            start_word = i * 250
            end_word = min((i + 1) * 250, len(words))
            page_text = " ".join(words[start_word:end_word])

            pages.append(
                {
                    "page_number": i + 1,
                    "text_length": len(page_text),
                    "estimated_words": len(page_text.split()),
                    "estimated": True,
                }
            )

        return pages

    def _calculate_docx_quality_metrics(
        self, text: str, metadata: Dict[str, Any]
    ) -> Dict[str, float]:
        """Calculate quality metrics for DOCX extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "word_count": len(text.split()),
            "paragraph_count": text.count("\n\n") + 1,
            "metadata_completeness": 0.0,
            "structure_quality": 0.0,
        }

        # Calculate metadata completeness
        metadata_fields = ["title", "author", "created"]
        filled_fields = sum(1 for field in metadata_fields if metadata.get(field))
        metrics["metadata_completeness"] = filled_fields / len(metadata_fields)

        # Calculate structure quality
        structure_indicators = [". ", "\n\n", ":", ";"]
        structure_score = sum(
            1 for indicator in structure_indicators if indicator in text
        ) / len(structure_indicators)
        metrics["structure_quality"] = min(1.0, structure_score)

        # Overall quality
        metrics["overall_quality"] = (
            min(1.0, metrics["word_count"] / 100) * 0.4  # Content volume
            + metrics["metadata_completeness"] * 0.3  # Metadata quality
            + metrics["structure_quality"] * 0.3  # Structure quality
        )

        return metrics


class XLSXExtractor(BaseExtractor):
    """High-precision XLSX text extraction with table structure analysis."""

    def supports_format(self, content_type: str) -> bool:
        """Check if XLSX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ]

    async def extract(
        self, file_data: bytes, options: ProcessingOptions
    ) -> ExtractionResult:
        """Extract text and structure from XLSX document."""
        start_time = datetime.now()

        try:
            workbook = load_workbook(io.BytesIO(file_data), data_only=True)

            content_parts = []
            metadata = {}
            structure = {}
            tables = []
            pages = []

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = self._extract_sheet_data(sheet, sheet_name)
                content_parts.append(sheet_data["text"])
                tables.extend(sheet_data["tables"])

            # Extract workbook properties
            metadata = self._extract_xlsx_metadata(workbook)

            # Combine all content
            full_text = "\n\n".join(filter(None, content_parts))

            # Create sheet information
            pages = self._create_xlsx_sheet_info(workbook)

            # Analyze workbook structure
            structure = self._analyze_xlsx_structure(workbook)

            # Calculate quality metrics
            quality_metrics = self._calculate_xlsx_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=pages,
                tables=tables,
                images=[],
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.95,
            )

        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise

    def _extract_sheet_data(self, sheet, sheet_name: str) -> Dict[str, Any]:
        """Extract data from a worksheet."""
        content_parts = []
        tables = []

        # Find data range
        data_range = self._find_data_range(sheet)
        if not data_range:
            return {"text": "", "tables": []}

        min_row, min_col, max_row, max_col = data_range

        # Convert to DataFrame for easier processing
        data = []
        for row in sheet.iter_rows(
            min_row=min_row,
            max_row=max_row,
            min_col=min_col,
            max_col=max_col,
            values_only=True,
        ):
            data.append(row)

        if data:
            # Create DataFrame
            df = pd.DataFrame(data[1:], columns=data[0] if len(data) > 1 else None)

            # Clean DataFrame
            df = df.dropna(how="all").dropna(axis=1, how="all")

            if not df.empty:
                # Convert to text representation
                table_text = self._dataframe_to_text(df)
                content_parts.append(f"Sheet: {sheet_name}\n{table_text}")

                # Create table structure
                table_data = {
                    "sheet_name": sheet_name,
                    "table_index": len(tables),
                    "headers": list(df.columns),
                    "rows": [list(row) for row in df.values],
                    "row_count": len(df),
                    "col_count": len(df.columns),
                    "text": table_text,
                    "confidence": 0.95,
                }
                tables.append(table_data)

        # Extract formulas and cell notes if present
        formulas_text = self._extract_formulas(sheet, data_range)
        if formulas_text:
            content_parts.append(f"Formulas in {sheet_name}:\n{formulas_text}")

        return {"text": "\n\n".join(filter(None, content_parts)), "tables": tables}

    def _find_data_range(self, sheet) -> Optional[Tuple[int, int, int, int]]:
        """Find the actual data range in a worksheet."""
        min_row, min_col = float("inf"), float("inf")
        max_row, max_col = 0, 0
        has_data = False

        for row in sheet.iter_rows():
            for cell in row:
                if cell.value is not None and str(cell.value).strip():
                    has_data = True
                    min_row = min(min_row, cell.row)
                    max_row = max(max_row, cell.row)
                    min_col = min(min_col, cell.column)
                    max_col = max(max_col, cell.column)

        if has_data:
            return (min_row, min_col, max_row, max_col)
        return None

    def _dataframe_to_text(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to readable text format."""
        text_parts = []

        # Add headers
        headers = [str(col) for col in df.columns]
        text_parts.append(" | ".join(headers))
        text_parts.append("-" * len(" | ".join(headers)))

        # Add data rows
        for _, row in df.iter_rows():
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            text_parts.append(row_text)

        return "\n".join(text_parts)

    def _extract_formulas(self, sheet, data_range: Tuple[int, int, int, int]) -> str:
        """Extract formulas from worksheet."""
        min_row, min_col, max_row, max_col = data_range
        formulas = []

        for row in sheet.iter_rows(
            min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col
        ):
            for cell in row:
                if cell.data_type == "f":  # Formula cell
                    cell_ref = f"{get_column_letter(cell.column)}{cell.row}"
                    formulas.append(f"{cell_ref}: {cell.formula} = {cell.value}")

        return "\n".join(formulas) if formulas else ""

    def _extract_xlsx_metadata(self, workbook) -> Dict[str, Any]:
        """Extract metadata from XLSX workbook."""
        metadata = {}

        try:
            props = workbook.properties
            metadata.update(
                {
                    "title": props.title or "",
                    "author": props.creator or "",
                    "subject": props.subject or "",
                    "description": props.description or "",
                    "keywords": props.keywords or "",
                    "category": props.category or "",
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                    "last_modified_by": props.lastModifiedBy or "",
                }
            )
        except Exception as e:
            logger.warning(f"Failed to extract workbook properties: {e}")

        # Add workbook statistics
        metadata.update(
            {
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
            }
        )

        return metadata

    def _create_xlsx_sheet_info(self, workbook) -> List[Dict[str, Any]]:
        """Create sheet information for XLSX."""
        sheets_info = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # Count non-empty cells
            non_empty_cells = 0
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        non_empty_cells += 1

            sheets_info.append(
                {
                    "sheet_name": sheet_name,
                    "page_number": len(sheets_info) + 1,
                    "non_empty_cells": non_empty_cells,
                    "max_row": sheet.max_row or 0,
                    "max_column": sheet.max_column or 0,
                    "has_data": non_empty_cells > 0,
                }
            )

        return sheets_info

    def _analyze_xlsx_structure(self, workbook) -> Dict[str, Any]:
        """Analyze XLSX workbook structure."""
        structure = {
            "has_formulas": False,
            "has_charts": False,
            "has_pivot_tables": False,
            "sheet_types": {},
            "data_density": 0.0,
        }

        total_cells = 0
        data_cells = 0

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]

            # Count cells
            sheet_total = 0
            sheet_data = 0
            has_formulas = False

            for row in sheet.iter_rows():
                for cell in row:
                    sheet_total += 1
                    if cell.value is not None:
                        sheet_data += 1
                    if cell.data_type == "f":
                        has_formulas = True

            total_cells += sheet_total
            data_cells += sheet_data

            if has_formulas:
                structure["has_formulas"] = True

            # Determine sheet type
            if sheet_data == 0:
                sheet_type = "empty"
            elif has_formulas:
                sheet_type = "calculation"
            elif sheet_data / max(sheet_total, 1) > 0.5:
                sheet_type = "data"
            else:
                sheet_type = "sparse"

            structure["sheet_types"][sheet_name] = sheet_type

        # Calculate overall data density
        structure["data_density"] = data_cells / max(total_cells, 1)

        return structure

    def _calculate_xlsx_quality_metrics(
        self, text: str, metadata: Dict[str, Any]
    ) -> Dict[str, float]:
        """Calculate quality metrics for XLSX extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "cell_count": text.count("|") // 2,  # Rough estimate
            "sheet_count": metadata.get("sheet_count", 0),
            "data_density": 0.0,
            "structure_quality": 0.0,
        }

        # Calculate data density from structure analysis
        if "structure" in metadata:
            metrics["data_density"] = metadata["structure"].get("data_density", 0.0)

        # Calculate structure quality
        structure_indicators = ["Sheet:", "|", "-"]  # Table-like structure
        structure_score = sum(
            1 for indicator in structure_indicators if indicator in text
        ) / len(structure_indicators)
        metrics["structure_quality"] = min(1.0, structure_score)

        # Overall quality
        metrics["overall_quality"] = (
            min(1.0, metrics["cell_count"] / 100) * 0.4  # Data volume
            + metrics["data_density"] * 0.3  # Data density
            + metrics["structure_quality"] * 0.3  # Structure quality
        )

        return metrics


class PPTXExtractor(BaseExtractor):
    """High-precision PPTX text extraction with slide structure analysis."""

    def supports_format(self, content_type: str) -> bool:
        """Check if PPTX format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]

    async def extract(
        self, file_data: bytes, options: ProcessingOptions
    ) -> ExtractionResult:
        """Extract text and structure from PPTX document."""
        start_time = datetime.now()

        try:
            presentation = Presentation(io.BytesIO(file_data))

            content_parts = []
            metadata = {}
            structure = {}
            tables = []
            pages = []

            # Process each slide
            for slide_idx, slide in enumerate(presentation.slides):
                slide_data = self._extract_slide_data(slide, slide_idx)
                content_parts.append(slide_data["text"])
                tables.extend(slide_data["tables"])

            # Extract presentation properties
            metadata = self._extract_pptx_metadata(presentation)

            # Combine all content
            full_text = "\n\n".join(filter(None, content_parts))

            # Create slide information
            pages = self._create_pptx_slide_info(presentation)

            # Analyze presentation structure
            structure = self._analyze_pptx_structure(presentation)

            # Calculate quality metrics
            quality_metrics = self._calculate_pptx_quality_metrics(full_text, metadata)

            processing_time = int((datetime.now() - start_time).total_seconds() * 1000)

            return ExtractionResult(
                text=full_text,
                metadata=metadata,
                pages=pages,
                tables=tables,
                images=[],
                structure=structure,
                quality_metrics=quality_metrics,
                processing_time_ms=processing_time,
                confidence=0.90,
            )

        except PackageNotFoundError:
            logger.error("Invalid PPTX file format")
            raise ValueError("Invalid PPTX file format")
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise

    def _extract_slide_data(self, slide, slide_idx: int) -> Dict[str, Any]:
        """Extract data from a slide."""
        content_parts = []
        tables = []

        # Add slide header
        content_parts.append(f"Slide {slide_idx + 1}")

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape_text = shape.text.strip()

                # Categorize text based on shape type and position
                if self._is_title_shape(shape):
                    content_parts.append(f"TITLE: {shape_text}")
                elif self._is_subtitle_shape(shape):
                    content_parts.append(f"SUBTITLE: {shape_text}")
                else:
                    content_parts.append(shape_text)

            # Extract tables
            if shape.has_table:
                table_data = self._extract_shape_table(shape, slide_idx)
                tables.append(table_data)
                content_parts.append(table_data["text"])

            # Extract notes (if any)
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                notes_text = self._extract_slide_notes(slide.notes_slide)
                if notes_text:
                    content_parts.append(f"NOTES: {notes_text}")

        return {"text": "\n".join(content_parts), "tables": tables}

    def _is_title_shape(self, shape) -> bool:
        """Check if shape is likely a title."""
        if hasattr(shape, "placeholder_format"):
            try:
                return shape.placeholder_format.type == 1  # Title placeholder
            except:
                pass

        # Fallback: check position and size
        if hasattr(shape, "top") and hasattr(shape, "height"):
            return shape.top < 100000  # Near top of slide

        return False

    def _is_subtitle_shape(self, shape) -> bool:
        """Check if shape is likely a subtitle."""
        if hasattr(shape, "placeholder_format"):
            try:
                return shape.placeholder_format.type == 2  # Subtitle placeholder
            except:
                pass
        return False

    def _extract_shape_table(self, shape, slide_idx: int) -> Dict[str, Any]:
        """Extract table from shape."""
        if not shape.has_table:
            return {}

        table = shape.table
        table_data = []

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_data.append(cell_text)
            table_data.append(row_data)

        # Convert to text
        table_text = []
        for row in table_data:
            table_text.append(" | ".join(row))

        return {
            "slide_number": slide_idx + 1,
            "table_index": 0,  # Usually only one table per slide
            "rows": table_data,
            "row_count": len(table_data),
            "col_count": len(table_data[0]) if table_data else 0,
            "text": "\n".join(table_text),
            "confidence": 0.90,
        }

    def _extract_slide_notes(self, notes_slide) -> str:
        """Extract text from slide notes."""
        notes_parts = []

        for shape in notes_slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                notes_parts.append(shape.text.strip())

        return " ".join(notes_parts)

    def _extract_pptx_metadata(self, presentation) -> Dict[str, Any]:
        """Extract metadata from PPTX presentation."""
        metadata = {}

        try:
            core_props = presentation.core_properties
            metadata.update(
                {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "keywords": core_props.keywords or "",
                    "category": core_props.category or "",
                    "comments": core_props.comments or "",
                    "created": core_props.created.isoformat()
                    if core_props.created
                    else None,
                    "modified": core_props.modified.isoformat()
                    if core_props.modified
                    else None,
                    "last_modified_by": core_props.last_modified_by or "",
                }
            )
        except Exception as e:
            logger.warning(f"Failed to extract presentation properties: {e}")

        # Add presentation statistics
        metadata.update(
            {
                "slide_count": len(presentation.slides),
            }
        )

        return metadata

    def _create_pptx_slide_info(self, presentation) -> List[Dict[str, Any]]:
        """Create slide information for PPTX."""
        slides_info = []

        for slide_idx, slide in enumerate(presentation.slides):
            # Count text elements
            text_elements = 0
            word_count = 0

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_elements += 1
                    word_count += len(shape.text.split())

            slides_info.append(
                {
                    "slide_number": slide_idx + 1,
                    "text_length": word_count * 5,  # Rough estimate
                    "text_elements": text_elements,
                    "word_count": word_count,
                    "has_content": text_elements > 0,
                }
            )

        return slides_info

    def _analyze_pptx_structure(self, presentation) -> Dict[str, Any]:
        """Analyze PPTX presentation structure."""
        structure = {
            "has_title_slide": False,
            "has_content_slides": False,
            "has_conclusion_slide": False,
            "slide_types": {},
            "avg_words_per_slide": 0.0,
            "presentation_type": "unknown",
        }

        total_words = 0
        content_slides = 0

        for slide_idx, slide in enumerate(presentation.slides):
            slide_words = 0
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    slide_words += len(text.split())

            total_words += slide_words

            # Determine slide type
            if slide_idx == 0 and slide_words < 20:
                structure["has_title_slide"] = True
                slide_type = "title"
            elif any(
                keyword in " ".join(slide_text).lower()
                for keyword in ["conclusion", "summary", "thank", "questions"]
            ):
                structure["has_conclusion_slide"] = True
                slide_type = "conclusion"
            elif slide_words > 10:
                structure["has_content_slides"] = True
                content_slides += 1
                slide_type = "content"
            else:
                slide_type = "other"

            structure["slide_types"][f"slide_{slide_idx + 1}"] = slide_type

        # Calculate averages
        structure["avg_words_per_slide"] = total_words / max(
            len(presentation.slides), 1
        )

        # Determine presentation type
        if structure["has_title_slide"] and structure["has_content_slides"]:
            if structure["avg_words_per_slide"] < 50:
                structure["presentation_type"] = "visual"
            else:
                structure["presentation_type"] = "detailed"
        else:
            structure["presentation_type"] = "simple"

        return structure

    def _calculate_pptx_quality_metrics(
        self, text: str, metadata: Dict[str, Any]
    ) -> Dict[str, float]:
        """Calculate quality metrics for PPTX extraction."""
        if not text:
            return {"overall_quality": 0.0}

        metrics = {
            "text_length": len(text),
            "word_count": len(text.split()),
            "slide_count": metadata.get("slide_count", 0),
            "structure_quality": 0.0,
            "content_completeness": 0.0,
        }

        # Calculate structure quality
        structure_indicators = ["Slide", "TITLE:", "SUBTITLE:", "NOTES:"]
        structure_score = sum(
            1 for indicator in structure_indicators if indicator in text
        ) / len(structure_indicators)
        metrics["structure_quality"] = min(1.0, structure_score)

        # Calculate content completeness
        if metrics["slide_count"] > 0:
            avg_words_per_slide = metrics["word_count"] / metrics["slide_count"]
            metrics["content_completeness"] = min(
                1.0, avg_words_per_slide / 50
            )  # 50 words per slide is good
        else:
            metrics["content_completeness"] = 0.0

        # Overall quality
        metrics["overall_quality"] = (
            min(1.0, metrics["word_count"] / 100) * 0.4  # Content volume
            + metrics["structure_quality"] * 0.3  # Structure quality
            + metrics["content_completeness"] * 0.3  # Content completeness
        )

        return metrics
