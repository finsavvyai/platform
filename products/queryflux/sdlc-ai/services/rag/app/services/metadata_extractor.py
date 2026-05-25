"""
Comprehensive metadata preservation and extraction service for SDLC.ai platform.

This module provides extensive metadata extraction capabilities for various document types,
preserving structural information, document properties, and content relationships.
"""

import hashlib
import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple, Union

import magic
from PIL import Image
from PIL.ExifTags import TAGS

logger = logging.getLogger(__name__)


@dataclass
class DocumentMetadata:
    """Comprehensive document metadata structure."""

    # Basic file information
    filename: str
    file_size: int
    content_type: str
    encoding: Optional[str] = None
    checksum_md5: str = ""
    checksum_sha256: str = ""

    # Document properties
    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    keywords: Optional[str] = None
    description: Optional[str] = None
    category: Optional[str] = None
    comments: Optional[str] = None

    # Creation and modification information
    created_date: Optional[datetime] = None
    modified_date: Optional[datetime] = None
    last_modified_by: Optional[str] = None
    revision_number: Optional[str] = None

    # Technical properties
    language: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None
    paragraph_count: Optional[int] = None

    # Structure information
    has_tables: bool = False
    has_images: bool = False
    has_headers: bool = False
    has_footers: bool = False
    has_toc: bool = False
    has_bookmarks: bool = False
    has_annotations: bool = False
    has_forms: bool = False

    # Content analysis
    content_density: float = 0.0
    readability_score: float = 0.0
    complexity_score: float = 0.0
    technical_terms: List[str] = field(default_factory=list)
    named_entities: List[Dict[str, Any]] = field(default_factory=list)

    # Security and classification
    security_classification: Optional[str] = None
    access_level: Optional[str] = None
    retention_period: Optional[str] = None
    compliance_tags: List[str] = field(default_factory=list)

    # Processing information
    extraction_method: str = ""
    extraction_confidence: float = 1.0
    processing_time_ms: int = 0
    processing_errors: List[str] = field(default_factory=list)

    # Custom metadata
    custom_properties: Dict[str, Any] = field(default_factory=dict)
    raw_metadata: Dict[str, Any] = field(default_factory=dict)


class BaseMetadataExtractor(ABC):
    """Abstract base class for metadata extractors."""

    @abstractmethod
    def supports_format(self, content_type: str) -> bool:
        """Check if the extractor supports the given content type."""
        pass

    @abstractmethod
    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract metadata from file data."""
        pass


class CommonMetadataExtractor(BaseMetadataExtractor):
    """Extract common metadata applicable to all file types."""

    def supports_format(self, content_type: str) -> bool:
        """Supports all file types as a fallback."""
        return True

    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract common metadata from file data."""
        metadata = DocumentMetadata(
            filename=filename,
            file_size=len(file_data),
            content_type=self._detect_content_type(file_data),
        )

        # Calculate checksums
        metadata.checksum_md5 = hashlib.md5(file_data).hexdigest()
        metadata.checksum_sha256 = hashlib.sha256(file_data).hexdigest()

        # Detect encoding for text files
        if metadata.content_type.startswith("text/"):
            metadata.encoding = self._detect_encoding(file_data)

        # Extract general file properties
        metadata.custom_properties.update(
            {
                "file_extension": self._get_file_extension(filename),
                "mime_type": metadata.content_type,
                "is_binary": not self._is_text_file(file_data),
            }
        )

        return metadata

    def _detect_content_type(self, file_data: bytes) -> str:
        """Detect content type from file data."""
        try:
            mime = magic.Magic(mime=True)
            return mime.from_buffer(file_data)
        except:
            # Fallback to basic detection
            if file_data.startswith(b"%PDF"):
                return "application/pdf"
            elif file_data.startswith(b"PK\x03\x04"):
                return "application/zip"  # Could be DOCX, XLSX, PPTX
            elif self._is_text_file(file_data):
                return "text/plain"
            else:
                return "application/octet-stream"

    def _detect_encoding(self, file_data: bytes) -> str:
        """Detect text encoding."""
        try:
            import chardet

            result = chardet.detect(file_data[:10000])  # Sample first 10KB
            return result.get("encoding", "utf-8")
        except:
            return "utf-8"

    def _is_text_file(self, file_data: bytes) -> bool:
        """Check if file is likely a text file."""
        try:
            # Check for null bytes (common in binary files)
            if b"\x00" in file_data[:1000]:
                return False

            # Try to decode as UTF-8
            file_data[:1000].decode("utf-8")
            return True
        except:
            return False

    def _get_file_extension(self, filename: str) -> str:
        """Get file extension from filename."""
        return filename.split(".")[-1].lower() if "." in filename else ""


class PDFMetadataExtractor(BaseMetadataExtractor):
    """Extract metadata from PDF documents."""

    def supports_format(self, content_type: str) -> bool:
        """Check if PDF format is supported."""
        return content_type.lower() in [
            "application/pdf",
            "application/x-pdf",
        ]

    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract metadata from PDF document."""
        from pypdf import PdfReader

        metadata = DocumentMetadata(
            filename=filename,
            file_size=len(file_data),
            content_type="application/pdf",
            extraction_method="pypdf",
        )

        try:
            pdf_reader = PdfReader(io.BytesIO(file_data))

            # Extract document properties
            if pdf_reader.metadata:
                metadata.title = pdf_reader.metadata.get("/Title", "").strip() or None
                metadata.author = pdf_reader.metadata.get("/Author", "").strip() or None
                metadata.subject = (
                    pdf_reader.metadata.get("/Subject", "").strip() or None
                )
                metadata.keywords = (
                    pdf_reader.metadata.get("/Keywords", "").strip() or None
                )
                metadata.creator = (
                    pdf_reader.metadata.get("/Creator", "").strip() or None
                )
                metadata.producer = (
                    pdf_reader.metadata.get("/Producer", "").strip() or None
                )

                # Parse dates
                creation_date = pdf_reader.metadata.get("/CreationDate")
                if creation_date:
                    metadata.created_date = self._parse_pdf_date(creation_date)

                modification_date = pdf_reader.metadata.get("/ModDate")
                if modification_date:
                    metadata.modified_date = self._parse_pdf_date(modification_date)

            # Extract document statistics
            metadata.page_count = len(pdf_reader.pages)
            metadata.has_forms = any("/AcroForm" in page for page in pdf_reader.pages)
            metadata.is_encrypted = pdf_reader.is_encrypted

            # Analyze PDF structure
            structure_info = await self._analyze_pdf_structure(pdf_reader)
            metadata.has_tables = structure_info.get("has_tables", False)
            metadata.has_images = structure_info.get("has_images", False)
            metadata.has_bookmarks = structure_info.get("has_bookmarks", False)

            # Store raw metadata
            metadata.raw_metadata = {
                "pdf_version": getattr(pdf_reader, "pdf_header", {}).get(
                    "version", "unknown"
                ),
                "is_encrypted": pdf_reader.is_encrypted,
                "permissions": getattr(pdf_reader, "permissions", {}),
            }

        except Exception as e:
            logger.error(f"Failed to extract PDF metadata: {e}")
            metadata.processing_errors.append(
                f"PDF metadata extraction failed: {str(e)}"
            )

        return metadata

    def _parse_pdf_date(self, date_str: str) -> Optional[datetime]:
        """Parse PDF date string."""
        try:
            # PDF dates are in format: D:YYYYMMDDHHmmSSOHH'mm'
            # Simplified parsing
            if date_str.startswith("D:"):
                date_str = date_str[2:]

            # Extract basic date components
            if len(date_str) >= 8:
                year = int(date_str[:4])
                month = int(date_str[4:6])
                day = int(date_str[6:8])
                return datetime(year, month, day)
        except:
            pass
        return None

    async def _analyze_pdf_structure(self, pdf_reader) -> Dict[str, Any]:
        """Analyze PDF document structure."""
        structure = {
            "has_tables": False,
            "has_images": False,
            "has_bookmarks": False,
        }

        try:
            # Check for images
            for page in pdf_reader.pages[:5]:  # Check first 5 pages
                if "/Resources" in page and "/XObject" in page.get("/Resources", {}):
                    structure["has_images"] = True
                    break

            # Check for bookmarks (outlines)
            if hasattr(pdf_reader, "outline") and pdf_reader.outline:
                structure["has_bookmarks"] = True

            # Check for tables (heuristic)
            sample_text = ""
            for page in pdf_reader.pages[:3]:
                try:
                    sample_text += page.extract_text() + " "
                except:
                    continue

            # Look for table-like patterns
            table_indicators = ["|", "\t", "    ", "────"]
            structure["has_tables"] = any(
                indicator in sample_text for indicator in table_indicators
            )

        except Exception as e:
            logger.warning(f"Failed to analyze PDF structure: {e}")

        return structure


class ImageMetadataExtractor(BaseMetadataExtractor):
    """Extract metadata from image files."""

    def supports_format(self, content_type: str) -> bool:
        """Check if image format is supported."""
        return content_type.lower().startswith("image/")

    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract metadata from image file."""
        metadata = DocumentMetadata(
            filename=filename,
            file_size=len(file_data),
            content_type=self._detect_image_type(file_data),
            extraction_method="PIL",
        )

        try:
            with Image.open(io.BytesIO(file_data)) as img:
                # Basic image properties
                metadata.custom_properties.update(
                    {
                        "width": img.width,
                        "height": img.height,
                        "format": img.format,
                        "mode": img.mode,
                        "has_transparency": img.mode in ("RGBA", "LA")
                        or "transparency" in img.info,
                    }
                )

                # Extract EXIF data
                exif_data = img._getexif()
                if exif_data:
                    exif_metadata = {}
                    for tag_id, value in exif_data.items():
                        tag = TAGS.get(tag_id, tag_id)
                        exif_metadata[tag] = value

                    # Map common EXIF fields to our metadata structure
                    metadata.title = exif_metadata.get("ImageDescription")
                    metadata.created_date = self._parse_exif_date(
                        exif_metadata.get("DateTime")
                    )
                    metadata.custom_properties["camera_make"] = exif_metadata.get(
                        "Make"
                    )
                    metadata.custom_properties["camera_model"] = exif_metadata.get(
                        "Model"
                    )
                    metadata.custom_properties["software"] = exif_metadata.get(
                        "Software"
                    )

                    metadata.raw_metadata["exif"] = exif_metadata

        except Exception as e:
            logger.error(f"Failed to extract image metadata: {e}")
            metadata.processing_errors.append(
                f"Image metadata extraction failed: {str(e)}"
            )

        return metadata

    def _detect_image_type(self, file_data: bytes) -> str:
        """Detect specific image type."""
        try:
            with Image.open(io.BytesIO(file_data)) as img:
                return f"image/{img.format.lower()}"
        except:
            return "image/unknown"

    def _parse_exif_date(self, date_str: Optional[str]) -> Optional[datetime]:
        """Parse EXIF date string."""
        if not date_str:
            return None

        try:
            # EXIF dates are in format: YYYY:MM:DD HH:MM:SS
            return datetime.strptime(date_str, "%Y:%m:%d %H:%M:%S")
        except:
            return None


class OfficeMetadataExtractor(BaseMetadataExtractor):
    """Extract metadata from Microsoft Office documents."""

    def supports_format(self, content_type: str) -> bool:
        """Check if Office format is supported."""
        return content_type.lower() in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/msword",
            "application/vnd.ms-excel",
            "application/vnd.ms-powerpoint",
        ]

    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract metadata from Office document."""
        metadata = DocumentMetadata(
            filename=filename,
            file_size=len(file_data),
            content_type=self._detect_office_type(file_data, filename),
            extraction_method="office_openxml",
        )

        try:
            if metadata.content_type.endswith("wordprocessingml.document"):
                await self._extract_docx_metadata(file_data, metadata)
            elif metadata.content_type.endswith("spreadsheetml.sheet"):
                await self._extract_xlsx_metadata(file_data, metadata)
            elif metadata.content_type.endswith("presentationml.presentation"):
                await self._extract_pptx_metadata(file_data, metadata)
            else:
                # Try legacy formats
                await self._extract_legacy_office_metadata(file_data, metadata)

        except Exception as e:
            logger.error(f"Failed to extract Office metadata: {e}")
            metadata.processing_errors.append(
                f"Office metadata extraction failed: {str(e)}"
            )

        return metadata

    def _detect_office_type(self, file_data: bytes, filename: str) -> str:
        """Detect specific Office document type."""
        ext = filename.lower().split(".")[-1]

        type_mapping = {
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "doc": "application/msword",
            "xls": "application/vnd.ms-excel",
            "ppt": "application/vnd.ms-powerpoint",
        }

        return type_mapping.get(ext, "application/octet-stream")

    async def _extract_docx_metadata(
        self, file_data: bytes, metadata: DocumentMetadata
    ):
        """Extract DOCX metadata."""
        from docx import Document

        doc = Document(io.BytesIO(file_data))

        # Extract core properties
        if doc.core_properties:
            metadata.title = doc.core_properties.title or None
            metadata.author = doc.core_properties.author or None
            metadata.subject = doc.core_properties.subject or None
            metadata.keywords = doc.core_properties.keywords or None
            metadata.category = doc.core_properties.category or None
            metadata.comments = doc.core_properties.comments or None
            metadata.created_date = doc.core_properties.created
            metadata.modified_date = doc.core_properties.modified
            metadata.last_modified_by = doc.core_properties.last_modified_by
            metadata.revision_number = (
                str(doc.core_properties.revision)
                if doc.core_properties.revision
                else None
            )

        # Document statistics
        metadata.paragraph_count = len(doc.paragraphs)
        metadata.word_count = sum(
            len(p.text.split()) for p in doc.paragraphs if p.text.strip()
        )
        metadata.character_count = sum(
            len(p.text) for p in doc.paragraphs if p.text.strip()
        )

        # Structure analysis
        metadata.has_tables = len(doc.tables) > 0
        metadata.has_headers = any(
            "heading" in p.style.name.lower() for p in doc.paragraphs if p.style
        )

        # Custom properties
        metadata.custom_properties.update(
            {
                "section_count": len(doc.sections),
                "table_count": len(doc.tables),
            }
        )

    async def _extract_xlsx_metadata(
        self, file_data: bytes, metadata: DocumentMetadata
    ):
        """Extract XLSX metadata."""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(file_data), read_only=True)

        # Extract core properties
        if wb.properties:
            metadata.title = wb.properties.title or None
            metadata.author = wb.properties.creator or None
            metadata.subject = wb.properties.subject or None
            metadata.description = wb.properties.description or None
            metadata.category = wb.properties.category or None
            metadata.created_date = wb.properties.created
            metadata.modified_date = wb.properties.modified
            metadata.last_modified_by = wb.properties.lastModifiedBy

        # Workbook statistics
        metadata.custom_properties.update(
            {
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
            }
        )

        # Analyze data content
        total_cells = 0
        data_cells = 0

        for sheet_name in wb.sheetnames[:3]:  # Analyze first 3 sheets
            sheet = wb[sheet_name]
            for row in sheet.iter_rows():
                for cell in row:
                    total_cells += 1
                    if cell.value is not None:
                        data_cells += 1

        if total_cells > 0:
            metadata.content_density = data_cells / total_cells

    async def _extract_pptx_metadata(
        self, file_data: bytes, metadata: DocumentMetadata
    ):
        """Extract PPTX metadata."""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_data))

        # Extract core properties
        if prs.core_properties:
            metadata.title = prs.core_properties.title or None
            metadata.author = prs.core_properties.author or None
            metadata.subject = prs.core_properties.subject or None
            metadata.keywords = prs.core_properties.keywords or None
            metadata.category = prs.core_properties.category or None
            metadata.comments = prs.core_properties.comments or None
            metadata.created_date = prs.core_properties.created
            metadata.modified_date = prs.core_properties.modified
            metadata.last_modified_by = prs.core_properties.last_modifiedBy

        # Presentation statistics
        slide_count = len(prs.slides)
        metadata.custom_properties.update(
            {
                "slide_count": slide_count,
            }
        )

        # Analyze slide content
        total_text_elements = 0
        has_tables = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    total_text_elements += 1
                if shape.has_table:
                    has_tables = True

        metadata.has_tables = has_tables
        metadata.word_count = total_text_elements * 5  # Rough estimate

    async def _extract_legacy_office_metadata(
        self, file_data: bytes, metadata: DocumentMetadata
    ):
        """Extract metadata from legacy Office formats."""
        # For legacy formats, we can extract basic information
        # Full extraction would require additional libraries like oletools
        metadata.extraction_method = "legacy_fallback"
        metadata.processing_errors.append(
            "Legacy Office format detected - limited metadata extraction"
        )


class TextMetadataExtractor(BaseMetadataExtractor):
    """Extract metadata from text files."""

    def supports_format(self, content_type: str) -> bool:
        """Check if text format is supported."""
        return content_type.lower().startswith("text/")

    async def extract(self, file_data: bytes, filename: str) -> DocumentMetadata:
        """Extract metadata from text file."""
        metadata = DocumentMetadata(
            filename=filename,
            file_size=len(file_data),
            content_type=self._detect_text_type(file_data, filename),
            extraction_method="text_analysis",
        )

        try:
            # Decode text
            encoding = self._detect_encoding(file_data)
            text = file_data.decode(encoding, errors="ignore")

            # Basic statistics
            metadata.character_count = len(text)
            metadata.word_count = len(text.split())
            metadata.paragraph_count = len([p for p in text.split("\n\n") if p.strip()])
            metadata.encoding = encoding

            # Language detection
            try:
                from langdetect import detect

                metadata.language = detect(text[:1000])  # Use first 1000 characters
            except:
                metadata.language = "unknown"

            # Content analysis
            metadata.content_density = self._calculate_text_density(text)
            metadata.readability_score = self._calculate_readability(text)
            metadata.complexity_score = self._calculate_complexity(text)

            # Extract structure information
            structure = self._analyze_text_structure(text)
            metadata.has_tables = structure.get("has_tables", False)
            metadata.has_headers = structure.get("has_headers", False)

            # Extract potential metadata from content
            content_metadata = self._extract_content_metadata(text)
            if content_metadata.get("title"):
                metadata.title = content_metadata["title"]
            if content_metadata.get("author"):
                metadata.author = content_metadata["author"]

        except Exception as e:
            logger.error(f"Failed to extract text metadata: {e}")
            metadata.processing_errors.append(
                f"Text metadata extraction failed: {str(e)}"
            )

        return metadata

    def _detect_text_type(self, file_data: bytes, filename: str) -> str:
        """Detect specific text file type."""
        ext = filename.lower().split(".")[-1]

        type_mapping = {
            "txt": "text/plain",
            "md": "text/markdown",
            "json": "application/json",
            "xml": "text/xml",
            "html": "text/html",
            "htm": "text/html",
            "csv": "text/csv",
            "log": "text/plain",
        }

        return type_mapping.get(ext, "text/plain")

    def _detect_encoding(self, file_data: bytes) -> str:
        """Detect text encoding."""
        try:
            import chardet

            result = chardet.detect(file_data[:10000])
            return result.get("encoding", "utf-8")
        except:
            return "utf-8"

    def _calculate_text_density(self, text: str) -> float:
        """Calculate text density (meaningful content ratio)."""
        if not text:
            return 0.0

        meaningful_chars = sum(1 for c in text if c.isalnum() or c.isspace())
        return meaningful_chars / len(text)

    def _calculate_readability(self, text: str) -> float:
        """Calculate basic readability score."""
        try:
            import textstat

            # Normalize to 0-1 range
            flesch_score = textstat.flesch_reading_ease(text)
            return min(1.0, max(0.0, flesch_score / 100))
        except:
            return 0.5  # Default score

    def _calculate_complexity(self, text: str) -> float:
        """Calculate text complexity based on vocabulary and sentence structure."""
        words = text.lower().split()
        if not words:
            return 0.0

        # Vocabulary diversity
        unique_words = len(set(words))
        vocab_diversity = unique_words / len(words)

        # Average word length
        avg_word_length = sum(len(word) for word in words) / len(words)
        length_score = min(1.0, avg_word_length / 10)  # Normalize

        # Sentence complexity
        sentences = text.split(". ")
        avg_sentence_length = (
            sum(len(s.split()) for s in sentences) / len(sentences) if sentences else 0
        )
        sentence_score = min(1.0, avg_sentence_length / 20)  # Normalize

        return vocab_diversity * 0.4 + length_score * 0.3 + sentence_score * 0.3

    def _analyze_text_structure(self, text: str) -> Dict[str, Any]:
        """Analyze text structure."""
        structure = {
            "has_tables": False,
            "has_headers": False,
        }

        # Check for tables
        lines = text.split("\n")
        table_indicators = ["|", "\t", "    ", "────"]
        structure["has_tables"] = any(
            any(indicator in line for indicator in table_indicators) for line in lines
        )

        # Check for headers (markdown style)
        header_patterns = [r"^#{1,6}\s+", r"^[A-Z][A-Z\s]*$"]
        structure["has_headers"] = any(
            re.search(pattern, line, re.MULTILINE)
            for pattern in header_patterns
            for line in lines
        )

        return structure

    def _extract_content_metadata(self, text: str) -> Dict[str, Any]:
        """Extract metadata from text content."""
        metadata = {}

        lines = text.split("\n")

        # Look for title in first few lines
        for line in lines[:10]:
            line = line.strip()
            if len(line) > 5 and len(line) < 100:
                if not metadata.get("title") and not line.endswith((".", ",", ";")):
                    metadata["title"] = line

        # Look for author information
        author_patterns = [
            r"(?i)author[:\s]+(.+)",
            r"(?i)by[:\s]+(.+)",
            r"(?i)written\s+by[:\s]+(.+)",
        ]

        for pattern in author_patterns:
            for line in lines[:20]:
                match = re.search(pattern, line)
                if match:
                    metadata["author"] = match.group(1).strip()
                    break

        return metadata


class MetadataExtractionService:
    """Main metadata extraction service."""

    def __init__(self):
        self.extractors = [
            PDFMetadataExtractor(),
            ImageMetadataExtractor(),
            OfficeMetadataExtractor(),
            TextMetadataExtractor(),
            CommonMetadataExtractor(),  # Fallback
        ]

    async def extract_metadata(
        self, file_data: bytes, filename: str, content_type: Optional[str] = None
    ) -> DocumentMetadata:
        """Extract comprehensive metadata from a file."""
        if not content_type:
            # Detect content type if not provided
            try:
                content_type = magic.Magic(mime=True).from_buffer(file_data)
            except:
                content_type = "application/octet-stream"

        # Find appropriate extractor
        extractor = None
        for ext in self.extractors:
            if ext.supports_format(content_type):
                extractor = ext
                break

        if not extractor:
            logger.warning(
                f"No specific extractor found for {content_type}, using common extractor"
            )
            extractor = self.extractors[-1]  # Common extractor as fallback

        # Extract metadata
        start_time = datetime.now()
        metadata = await extractor.extract(file_data, filename)
        metadata.processing_time_ms = int(
            (datetime.now() - start_time).total_seconds() * 1000
        )

        # Add service metadata
        metadata.custom_properties.update(
            {
                "extraction_timestamp": datetime.now().isoformat(),
                "extractor_used": extractor.__class__.__name__,
            }
        )

        return metadata

    async def extract_batch_metadata(
        self, files: List[Tuple[bytes, str, Optional[str]]]
    ) -> List[DocumentMetadata]:
        """Extract metadata from multiple files."""
        results = []

        for file_data, filename, content_type in files:
            try:
                metadata = await self.extract_metadata(
                    file_data, filename, content_type
                )
                results.append(metadata)
            except Exception as e:
                logger.error(f"Failed to extract metadata from {filename}: {e}")
                # Create minimal metadata for failed files
                error_metadata = DocumentMetadata(
                    filename=filename,
                    file_size=len(file_data),
                    content_type=content_type or "application/octet-stream",
                    processing_errors=[str(e)],
                )
                results.append(error_metadata)

        return results
